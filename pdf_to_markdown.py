import os
from pdf2image import convert_from_path  # 用于将PDF转换为图片
import pytesseract  # OCR文字识别工具
import requests  # 用于发送HTTP请求
import json
import fitz  # PyMuPDF，用于处理PDF文件
from typing import List, Optional, Union, Tuple
from PIL import Image  # 图片处理
import re  # 正则表达式
from pathlib import Path  # 路径处理
import langdetect  # 用于语言检测
import docx  # 导入python-docx库
from PIL import ImageEnhance
import time  # 用于添加请求间隔

# 导入配置
from config import (
    API_KEY, TESSERACT_CMD, POPPLER_PATH, API_URL, 
    MODEL_NAME, OCR_LANG, TEMPERATURE, MAX_TOKENS,
    SUPPORTED_FORMATS, SUPPORTED_IMAGE_FORMATS, SUPPORTED_TEXT_FORMATS,
    OCR_CONFIG, API_RETRY_COUNT, API_RETRY_DELAY, API_RATE_LIMIT,
    TOP_K, TOP_P, FREQUENCY_PENALTY, API_REQUEST_TIMEOUT,
    MAX_CHUNK_SIZE, LANGUAGE_DISPLAY_NAMES
)
from utils import (
    check_file_exists, ensure_directory_exists,
    merge_markdown_chunks, clean_markdown_format
)

class PDFToMarkdown:
    """PDF/图片/文本转Markdown工具类"""
    
    def __init__(self, progress_callback=None):
        """初始化配置"""
        self.progress_callback = progress_callback
        self.api_url = API_URL
        self.headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {API_KEY}"
        }
        
        # 设置 Tesseract 路径
        if TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
        self._init_ocr_config()
        self._init_language_patterns()
        self.ocr_language = 'chi_sim'  # 默认简体中文
        self.need_translation = False   # 默认不需要翻译
    
    def _init_ocr_config(self):
        """初始化OCR配置"""
        if os.name == 'nt':  # Windows系统
            self.poppler_path = POPPLER_PATH  # 设置Poppler路径
    
    def _init_language_patterns(self):
        """初始化不同语言的清理模式"""
        self.language_patterns = {
            'zh': {  # 中文（简体和繁体）
                'space': r'([^\u4e00-\u9fff])\s+([^\u4e00-\u9fff])',
                'punctuation': r'[""'']+',
                'noise': r'[·・︰]',
                'level2': r'[^\u4e00-\u9fff\u3000-\u303f\uff00-\uffef\u0020-\u007f\n]'
            },
            'en': {  # 英文
                'space': r'\s+',
                'punctuation': r'[""'']+',
                'noise': r'[·・︰]',
                'level2': r'[^\x20-\x7f\n]'
            },
            'ja': {  # 日文
                'space': r'([^\u4e00-\u9fff\u3040-\u309f\u30a0-\u30ff])\s+([^\u4e00-\u9fff\u3040-\u309f\u30a0-\u30ff])',
                'punctuation': r'[""'']+',
                'noise': r'[·・︰]',
                'level2': r'[^\u4e00-\u9fff\u3040-\u309f\u30a0-\u30ff\u3000-\u303f\uff00-\uffef\u0020-\u007f\n]'
            },
            'default': {  # 其他语言
                'space': r'\s+',
                'punctuation': r'[""'']+',
                'noise': r'[·・︰]',
                'level2': r'[^\x20-\x7f\n]'
            }
        }
    
    def _is_scanned_pdf(self, pdf_path: str) -> bool:
        """
        判断是否为扫描版PDF
        通过尝试提取文本来判断：如果无法提取文本，则认为是扫描版
        """
        doc = fitz.open(pdf_path)
        for page in doc:
            text = page.get_text()
            if not text.strip():  # 如果页面没有可提取的文本
                return True
        return False
    
    def extract_text_from_pdf(self, pdf_path: str) -> str:
        """
        从PDF中直接提取文本（用于非扫描版PDF）
        使用PyMuPDF提取每页文本并合并
        """
        doc = fitz.open(pdf_path)
        text = []
        for page in doc:
            text.append(page.get_text())
        return "\n".join(text)
    
    def convert_pdf_to_images(self, pdf_path: str) -> List[Image.Image]:
        """
        将PDF转换为图片列表
        用于扫描版PDF的处理
        """
        if not check_file_exists(pdf_path):
            raise FileNotFoundError(f"文件不存在: {pdf_path}")
            
        if os.name == 'nt':
            return convert_from_path(pdf_path, poppler_path=self.poppler_path)
        return convert_from_path(pdf_path)
    
    def set_ocr_language(self, language: str):
        """设置OCR识别语言"""
        self.ocr_language = language
    
    def set_need_translation(self, need_translation: bool):
        """设置是否需要翻译"""
        self.need_translation = need_translation
    
    def process_image(self, image_path: str) -> str:
        """处理图片文件"""
        try:
            # 打开图片
            image = Image.open(image_path)
            
            # 图像预处理
            processed_image = self.preprocess_image(image)
            
            # 设置OCR参数
            custom_config = f'-l {self.ocr_language} --oem 1 --psm 3 ' + \
                           '--dpi 300 ' + \
                           '-c preserve_interword_spaces=1 ' + \
                           '-c tessedit_char_blacklist=|' + \
                           '-c textord_heavy_nr=1 ' + \
                           '-c textord_min_linesize=2.5'  # 保留词间空格，去除容易误识别的字符，改进数字识别，改进小字体识别，改进OCR参数配置
            # OCR识别
            text = pytesseract.image_to_string(
                processed_image,
                lang=self.ocr_language,
                config=custom_config
            )
            
            return text
            
        except Exception as e:
            print(f"处理图片时出错: {str(e)}")
            raise
    
    def extract_text_from_images(self, images: Union[List[Image.Image], List[str]]) -> str:
        """从图片中提取文本"""
        extracted_text = []
        total = len(images)
        
        # 获取语言特定的OCR配置
        custom_config = OCR_CONFIG.get(self.ocr_language, '--oem 1 --psm 3')
        
        for i, image in enumerate(images, 1):
            try:
                if isinstance(image, str):  # 如果是图片路径
                    text = self.process_image(image)
                else:  # 如果是PIL Image对象
                    # 预处理图像
                    image = image.convert('L')  # 转灰度
                    # 增强对比度
                    enhancer = ImageEnhance.Contrast(image)
                    image = enhancer.enhance(2.0)
                    # 增强锐度
                    enhancer = ImageEnhance.Sharpness(image)
                    image = enhancer.enhance(2.0)
                    # 放大图片
                    width, height = image.size
                    image = image.resize((width*2, height*2), Image.Resampling.LANCZOS)
                    
                    text = pytesseract.image_to_string(
                        image,
                        config=custom_config,
                        lang=self.ocr_language
                    )
                extracted_text.append(text)
                print(f'处理进度: {i}/{total}')
            except Exception as e:
                print(f"处理第{i}张图片时出错: {str(e)}")
                
        return '\n'.join(extracted_text)
    
    def detect_language(self, text: str) -> str:
        """
        检测文本主要语言
        返回语言代码：
        'zh_cn' - 简体中文
        'zh_tw' - 繁体中文
        'en' - 英文
        'ja' - 日文
        'de' - 德语
        'fr' - 法语
        'ar' - 阿拉伯语
        """
        try:
            # 根据OCR语言设置优先判断
            if self.ocr_language == 'chi_sim':
                return 'zh_cn'
            elif self.ocr_language == 'chi_tra':
                return 'zh_tw'
            elif self.ocr_language == 'eng':
                return 'en'
            elif self.ocr_language == 'jpn':
                return 'ja'
            elif self.ocr_language == 'deu':
                return 'de'
            elif self.ocr_language == 'fra':
                return 'fr'
            elif self.ocr_language == 'ara':
                return 'ar'
            
            # 如果没有预设语言，则进行自动检测
            # 移除可能影响检测的内容
            cleaned_text = re.sub(r'[0-9\s\W]+', ' ', text)
            
            # 使用 langdetect 进行初步检测
            lang = langdetect.detect(cleaned_text)
            
            # 对中文进行进一步判断（简体/繁体）
            if lang == 'zh':
                # 计算简体和繁体字符的比例
                simplified = len([c for c in cleaned_text if '\u4e00' <= c <= '\u9fff' and c in simplified_chars])
                traditional = len([c for c in cleaned_text if '\u4e00' <= c <= '\u9fff' and c in traditional_chars])
                return 'zh_cn' if simplified >= traditional else 'zh_tw'
            
            # 映射语言代码
            lang_map = {
                'ja': 'ja',
                'de': 'de',
                'fr': 'fr',
                'ar': 'ar',
                'en': 'en'
            }
            
            return lang_map.get(lang, 'en')  # 默认返回英文
        except:
            # 如果检测失败，返回OCR设置的语言
            if self.ocr_language.startswith('chi_'):
                return 'zh_cn' if self.ocr_language == 'chi_sim' else 'zh_tw'
            return 'en'  # 其他情况默认返回英文
    
    def is_text_file(self, file_path: str) -> bool:
        """判断是否为文本文件（包括Word文档）"""
        return Path(file_path).suffix.lower() in SUPPORTED_TEXT_FORMATS
    
    def read_text_file(self, file_path: str) -> Tuple[str, str]:
        """
        读取文本文件并检测编码
        支持Word文档和普通文本文件
        返回：(文本内容, 检测到的语言)
        """
        file_ext = Path(file_path).suffix.lower()
        
        # 处理Word文档
        if file_ext in ['.doc', '.docx']:
            try:
                doc = docx.Document(file_path)
                content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                language = self.detect_language(content)
                return content, language
            except Exception as e:
                print(f"处理Word文档时出错: {str(e)}")
                raise
        
        # 处理其他文本文件
        encodings = ['utf-8', 'gbk', 'gb2312', 'iso-8859-1']
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                    language = self.detect_language(content)
                    return content, language
            except UnicodeDecodeError:
                continue
        
        raise ValueError(f"无法以支持的编码格式读取文件: {file_path}")
    
    def clean_text(self, text: str, clean_level: int = 1) -> str:
        """根据语言和清理级别清理文本"""
        if clean_level == 0:
            return text
        
        # 获取语言对应的清理规则
        lang_code = self.ocr_language[:2]  # 获取语言代码前两位
        patterns = self.language_patterns.get(
            lang_code, 
            self.language_patterns['default']
        )
        
        # 基础清理（所有语言通用）
        text = re.sub(r'\r\n', '\n', text)  # 统一换行符
        text = re.sub(r'\n{3,}', '\n\n', text)  # 合并多个空行
        
        # 根据语言特点清理
        if lang_code in ['zh', 'ja']:  # 中文和日文
            text = re.sub(patterns['space'], r'\1\2', text)  # 保留中日文间的空格
        else:  # 其他语言
            text = re.sub(patterns['space'], ' ', text)  # 合并空格
        
        # 清理标点符号
        text = re.sub(patterns['punctuation'], '"', text)
        
        if clean_level >= 1:
            # 清理确定的干扰字符
            text = re.sub(patterns['noise'], '', text)
            
        if clean_level >= 2:
            # 强化清理
            text = re.sub(patterns['level2'], '', text)
        
        return text.strip()
    
    def split_text_into_chunks(self, text: str, max_chunk_size: int = 2000, overlap: int = 200) -> List[str]:
        """
        将长文本分割成带重叠的块进行处理
        
        Args:
            text: 要分割的文本
            max_chunk_size: 每个块的最大字符数
            overlap: 块之间的重叠字符数，用于保持上下文连贯
            
        Returns:
            分割后的文本块列表
        """
        # 首先按段落分割
        paragraphs = text.split('\n\n')
        chunks = []
        current_chunk = []
        current_size = 0
        last_context = ""  # 用于存储上一个块的结尾内容
        
        for para in paragraphs:
            para_size = len(para)
            
            # 如果段落超过最大块大小，需要进一步分割
            if para_size > max_chunk_size:
                # 如果当前块不为空，先保存当前块
                if current_chunk:
                    chunk_text = '\n\n'.join(current_chunk)
                    chunks.append(chunk_text)
                    last_context = chunk_text[-overlap:] if len(chunk_text) > overlap else chunk_text
                    current_chunk = []
                    current_size = 0
                
                # 按句子分割长段落
                sentences = re.split(r'([。！？])', para)
                temp_para = last_context  # 添加上下文
                
                for i in range(0, len(sentences), 2):
                    sentence = sentences[i]
                    if i + 1 < len(sentences):
                        sentence += sentences[i+1]  # 加回标点符号
                    
                    if len(temp_para) + len(sentence) < max_chunk_size:
                        temp_para += sentence
                    else:
                        if temp_para:
                            chunks.append(temp_para)
                            last_context = temp_para[-overlap:] if len(temp_para) > overlap else temp_para
                        temp_para = last_context + sentence
            
                if temp_para:
                    current_chunk.append(temp_para)
                    current_size = len(temp_para)
            else:
                # 检查添加当前段落是否会超出块大小
                if current_size + para_size > max_chunk_size:
                    chunk_text = '\n\n'.join(current_chunk)
                    chunks.append(chunk_text)
                    last_context = chunk_text[-overlap:] if len(chunk_text) > overlap else chunk_text
                    current_chunk = [last_context + para]
                    current_size = len(current_chunk[0])
                else:
                    current_chunk.append(para)
                    current_size += para_size
        
        # 处理最后一个块
        if current_chunk:
            chunks.append('\n\n'.join(current_chunk))
        
        return chunks
    
    def convert_to_markdown(self, text: str) -> str:
        """将文本转换为Markdown格式"""
        try:
            chunks = self.split_text_into_chunks(text)
            markdown_chunks = []
            previous_context = ""  # 存储前一个块的处理结果
            
            for i, chunk in enumerate(chunks):
                print(f"正在处理文本块 {i+1}/{len(chunks)} ({len(chunk)} 字符)...")
                
                # 将前一个块的结果作为上下文
                context_message = f"请继续保持前文的格式和结构。前文的结尾是:\n\n{previous_context}\n\n" if previous_context else ""
                
                payload = {
                    "model": MODEL_NAME,
                    "messages": [
                        {
                            "role": "system",
                            "content": "你是一个文本格式转换专家。请将输入的文本转换为结构良好的Markdown格式，保持原文的层级结构和重要信息。注意保持标题层级的连贯性。"
                        },
                        {
                            "role": "user",
                            "content": f"{context_message}请将以下文本转换为Markdown格式，保持原有的结构和格式：\n\n{chunk}"
                        }
                    ],
                    "temperature": TEMPERATURE,
                    "max_tokens": MAX_TOKENS,
                    "top_k": TOP_K,
                    "top_p": TOP_P,
                    "frequency_penalty": FREQUENCY_PENALTY,
                    "stream": False
                }
                
                success = False
                for retry in range(API_RETRY_COUNT):
                    try:
                        print(f"发送请求 (尝试 {retry + 1}/{API_RETRY_COUNT})...")
                        response = requests.post(
                            self.api_url,
                            headers=self.headers,
                            json=payload,
                            timeout=(30, API_REQUEST_TIMEOUT)
                        )
                        
                        if response.status_code == 200:
                            result = response.json()
                            if 'choices' in result and len(result['choices']) > 0:
                                markdown_text = result['choices'][0]['message']['content']
                                markdown_chunks.append(markdown_text)
                                # 保存当前处理结果的最后部分作为上下文
                                previous_context = markdown_text[-500:] if len(markdown_text) > 500 else markdown_text
                                success = True
                                print("文本块处理成功")
                                break
                        else:
                            error_msg = response.json()
                            print(f"API错误 (状态码: {response.status_code}): {error_msg}")
                            if retry == API_RETRY_COUNT - 1:
                                raise Exception(f"API请求失败: {error_msg}")
                            time.sleep(API_RETRY_DELAY)
                            
                    except Exception as e:
                        print(f"请求出错 ({retry + 1}/{API_RETRY_COUNT}): {str(e)}")
                        if retry == API_RETRY_COUNT - 1:
                            raise
                        time.sleep(API_RETRY_DELAY)
                
                if not success:
                    print("文本块处理失败，保留原始文本")
                    markdown_chunks.append(chunk)
                
                # 添加请求间隔
                time.sleep(1/API_RATE_LIMIT)
            
            # 合并所有处理后的块并清理格式
            if markdown_chunks:
                print("合并处理后的文本块...")
                merged_text = merge_markdown_chunks(markdown_chunks)
                return clean_markdown_format(merged_text)
            else:
                return text
            
        except Exception as e:
            print(f"调用 API 时发生错误: {str(e)}")
            return text
    
    def translate_to_chinese(self, text: str) -> str:
        """将文本翻译成中文"""
        try:
            # 准备API请求
            payload = {
                "model": MODEL_NAME,
                "messages": [
                    {
                        "role": "system",
                        "content": "你是一个专业的翻译专家。请将输入的文本翻译成中文，保持原有的格式和结构。"
                    },
                    {
                        "role": "user",
                        "content": f"请将以下文本翻译成中文，保持Markdown格式：\n\n{text}"
                    }
                ],
                "temperature": TEMPERATURE,
                "top_k": TOP_K,
                "top_p": TOP_P,
                "frequency_penalty": FREQUENCY_PENALTY,
                "max_tokens": MAX_TOKENS
            }
            
            # 添加重试机制
            for retry in range(API_RETRY_COUNT):
                try:
                    response = requests.post(
                        self.api_url, 
                        headers=self.headers, 
                        json=payload,
                        timeout=API_REQUEST_TIMEOUT
                    )
                    response.raise_for_status()
                    result = response.json()
                    return result['choices'][0]['message']['content']
                except requests.exceptions.RequestException as e:
                    if retry == API_RETRY_COUNT - 1:
                        print(f"翻译API请求失败 ({retry + 1}/{API_RETRY_COUNT}): {str(e)}")
                        raise
                    print(f"翻译API请求失败，正在重试 ({retry + 1}/{API_RETRY_COUNT})")
                    time.sleep(API_RETRY_DELAY)
                
                # 添加请求间隔
                time.sleep(1/API_RATE_LIMIT)
            
        except Exception as e:
            print(f"翻译时发生错误: {str(e)}")
            return text
    
    def report_progress(self, progress: float):
        """报告处理进度"""
        if self.progress_callback:
            self.progress_callback(progress)
            
    def process_file(self, file_path: str, output_path: str, clean_level: int = 1) -> dict:
        """
        处理文件并转换为Markdown格式
        
        Args:
            file_path: 输入文件路径
            output_path: 输出文件路径
            clean_level: 文本清理级别（0-2）
            
        Returns:
            包含处理结果的字典
        """
        try:
            # 获取文件扩展名
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # 检查文件是否存在
            if not check_file_exists(file_path):
                raise FileNotFoundError(f"文件不存在: {file_path}")
            
            # 更新进度：开始处理
            self.report_progress(0)
            print(f"开始处理文件: {file_path}")
            
            # 根据文件类型选择处理方法
            if file_ext == '.pdf':
                raw_text = self.process_pdf(file_path)
            elif file_ext in SUPPORTED_FORMATS['image']:
                raw_text = self.process_image(file_path)
            elif file_ext in SUPPORTED_FORMATS['document']:
                raw_text = self.process_document(file_path)
            elif file_ext in SUPPORTED_FORMATS['presentation']:
                raw_text = self.process_presentation(file_path)
            elif file_ext in SUPPORTED_FORMATS['spreadsheet']:
                raw_text = self.process_spreadsheet(file_path)
            elif file_ext in SUPPORTED_FORMATS['ebook']:
                raw_text = self.process_ebook(file_path)
            else:
                raise ValueError(f"不支持的文件格式: {file_ext}")
            
            # 保存原始文本
            raw_path = output_path.replace('.md', '_raw.txt')
            with open(raw_path, 'w', encoding='utf-8') as f:
                f.write(raw_text)
            print(f"保存清理前的原始文本到: {raw_path}")
            
            # 检测语言
            detected_language = self.detect_language(raw_text)
            display_language = LANGUAGE_DISPLAY_NAMES.get(detected_language, detected_language)
            print(f"使用{display_language}语言规则清理文本...")
            
            # 清理文本
            cleaned_text = self.clean_text(raw_text, clean_level)
            
            # 转换为Markdown
            markdown_text = self.convert_to_markdown(cleaned_text)
            
            # 保存Markdown结果
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(markdown_text)
            print(f"保存Markdown文件到: {output_path}")
            
            result = {
                'original': markdown_text,
                'language': display_language
            }
            
            # 如果需要翻译
            if self.need_translation and not detected_language.startswith(('zh_cn', 'zh_tw')):
                print("正在翻译文本...")
                translated_text = self.translate_to_chinese(markdown_text)
                translated_path = output_path.replace('.md', '_zh.md')
                with open(translated_path, 'w', encoding='utf-8') as f:
                    f.write(translated_text)
                print(f"保存翻译后的文件到: {translated_path}")
                result['translated'] = translated_text
            
            # 更新进度：完成
            self.report_progress(100)
            
            return result
            
        except Exception as e:
            print(f"处理文件时发生错误: {str(e)}")
            raise
            
    def process_presentation(self, file_path: str) -> dict:
        """处理演示文稿文件"""
        if file_path.endswith(('.ppt', '.pptx')):
            from pptx import Presentation
            prs = Presentation(file_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                text_content.append(f"\n## 幻灯片 {i+1}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                        
            return self.process_text('\n'.join(text_content))
            
    def process_spreadsheet(self, file_path: str) -> dict:
        """处理电子表格文件"""
        if file_path.endswith(('.xls', '.xlsx')):
            import pandas as pd
            df = pd.read_excel(file_path)
            return self.process_text(df.to_markdown())
        elif file_path.endswith('.csv'):
            import pandas as pd
            df = pd.read_csv(file_path)
            return self.process_text(df.to_markdown())
            
    def process_ebook(self, file_path: str) -> dict:
        """处理电子书文件"""
        if file_path.endswith('.epub'):
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup
            
            book = epub.read_epub(file_path)
            chapters = []
            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    chapters.append(soup.get_text())
                    
            return self.process_text('\n\n'.join(chapters))

    def process_pdf(self, pdf_path: str) -> str:
        """
        处理PDF文件，返回提取的文本
        
        Args:
            pdf_path: PDF文件路径
            
        Returns:
            提取的文本内容
        """
        try:
            if not check_file_exists(pdf_path):
                raise FileNotFoundError(f"PDF文件不存在: {pdf_path}")
            
            # 检查是否为扫描版PDF
            if self._is_scanned_pdf(pdf_path):
                print("检测到扫描版PDF，使用OCR处理...")
                # 将PDF转换为图片
                images = self.convert_pdf_to_images(pdf_path)
                
                # 处理每一页
                text_parts = []
                total_pages = len(images)
                
                for i, image in enumerate(images, 1):
                    print(f"正在处理第 {i}/{total_pages} 页...")
                    
                    # 检测是否为红头文件正文部分
                    if self._is_red_header_page(image) and i == 1:
                        print("检测到红头文件正文，跳过处理...")
                        continue
                    
                    # 检测是否为附件部分
                    is_attachment = self._is_attachment_page(image) or i > 1
                    
                    if is_attachment:
                        print(f"处理附件内容: 第 {i} 页")
                        # 图像预处理
                        image = self.preprocess_image(image)
                        # OCR识别
                        page_text = pytesseract.image_to_string(
                            image,
                            lang=self.ocr_language,
                            config=OCR_CONFIG.get(self.ocr_language, '')
                        )
                        text_parts.append(page_text)
                    else:
                        print(f"跳过非附件内容: 第 {i} 页")
                    
                    # 更新进度
                    self.report_progress(i * 100 / total_pages)
                
                return '\n\n'.join(text_parts)
            
            else:
                print("检测到可直接提取文本的PDF...")
                # 直接提取文本，但只提取附件部分
                return self._extract_attachment_text(pdf_path)
            
        except Exception as e:
            print(f"处理PDF文件时发生错误: {str(e)}")
            raise

    def _is_red_header_page(self, image: Image.Image) -> bool:
        """
        检测是否为红头文件页面
        
        Args:
            image: 页面图像
            
        Returns:
            是否为红头文件页面
        """
        try:
            # 转换为RGB模式
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            # 获取图像上部区域
            top_region = image.crop((0, 0, image.width, int(image.height * 0.15)))
            
            # 计算红色像素比例
            red_pixels = 0
            total_pixels = top_region.width * top_region.height
            
            for x in range(top_region.width):
                for y in range(top_region.height):
                    r, g, b = top_region.getpixel((x, y))
                    # 检测红色像素 (R值高，G和B值低)
                    if r > 150 and g < 100 and b < 100:
                        red_pixels += 1
            
            red_ratio = red_pixels / total_pixels
            
            # 如果红色像素比例超过阈值，判定为红头文件
            return red_ratio > 0.2
            
        except Exception as e:
            print(f"检测红头文件时出错: {str(e)}")
            return False

    def _is_attachment_page(self, image: Image.Image) -> bool:
        """
        检测是否为附件页面
        
        Args:
            image: 页面图像
            
        Returns:
            是否为附件页面
        """
        try:
            # 转换为灰度图
            if image.mode != 'L':
                image = image.convert('L')
            
            # 提取页面顶部区域进行OCR
            top_region = image.crop((0, 0, image.width, int(image.height * 0.2)))
            
            # OCR识别
            text = pytesseract.image_to_string(
                top_region,
                lang=self.ocr_language
            )
            
            # 检查是否包含"附件"关键词
            attachment_keywords = ['附件', '附录', 'attachment', 'appendix']
            return any(keyword in text for keyword in attachment_keywords)
            
        except Exception as e:
            print(f"检测附件页面时出错: {str(e)}")
            return False

    def _extract_attachment_text(self, pdf_path: str) -> str:
        """
        从PDF中提取附件部分的文本
        
        Args:
            pdf_path: PDF文件路径
            
        Returns:
            附件部分的文本
        """
        doc = fitz.open(pdf_path)
        text_parts = []
        in_attachment = False
        
        for page_num, page in enumerate(doc):
            text = page.get_text()
            
            # 检查是否为第一页
            if page_num == 0:
                # 检查是否为红头文件
                if self._is_red_header_text(text):
                    continue
            
            # 检查是否包含附件标记
            if not in_attachment:
                if self._has_attachment_marker(text):
                    in_attachment = True
                    # 提取附件标记后的内容
                    attachment_start = self._find_attachment_start(text)
                    if attachment_start >= 0:
                        text = text[attachment_start:]
            
            if in_attachment or page_num > 0:
                text_parts.append(text)
        
        return "\n\n".join(text_parts)

    def _is_red_header_text(self, text: str) -> bool:
        """
        检查文本是否为红头文件正文
        
        Args:
            text: 页面文本
            
        Returns:
            是否为红头文件正文
        """
        # 检查常见的红头文件特征
        red_header_markers = [
            '关于', '通知', '报告', '意见', '决定', '批复',
            '函', '令', '公告', '通报', '请示'
        ]
        
        first_lines = text.split('\n')[:5]  # 取前5行
        first_text = '\n'.join(first_lines)
        
        return any(marker in first_text for marker in red_header_markers)

    def _has_attachment_marker(self, text: str) -> bool:
        """
        检查文本是否包含附件标记
        
        Args:
            text: 页面文本
            
        Returns:
            是否包含附件标记
        """
        attachment_markers = ['附件', '附录', 'attachment', 'appendix']
        return any(marker in text.lower() for marker in attachment_markers)

    def _find_attachment_start(self, text: str) -> int:
        """
        查找附件内容的起始位置
        
        Args:
            text: 页面文本
            
        Returns:
            附件内容的起始位置，如果未找到则返回-1
        """
        attachment_patterns = [
            r'附\s*件\s*[：:]\s*',
            r'附\s*件\s*\d+\s*[：:]\s*',
            r'attachment\s*[：:]\s*',
            r'appendix\s*[：:]\s*'
        ]
        
        for pattern in attachment_patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.end()
        
        return -1

    def preprocess_image(self, image: Image.Image) -> Image.Image:
        """
        图像预处理，优化OCR效果
        
        Args:
            image: 输入图像
            
        Returns:
            处理后的图像
        """
        try:
            # 1. 转换为灰度图
            image = image.convert('L')
            
            # 2. 增强对比度和亮度（针对黄色背景优化）
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(2.5)  # 增加对比度
            
            enhancer = ImageEnhance.Brightness(image)
            image = enhancer.enhance(1.2)  # 调整亮度
            
            # 3. 自适应二值化（改进算法）
            import numpy as np
            import cv2
            
            # 转换为OpenCV格式
            img_array = np.array(image)
            
            # 使用OpenCV的自适应二值化
            img_array = cv2.adaptiveThreshold(
                img_array,
                255,
                cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                cv2.THRESH_BINARY,
                blockSize=15,  # 适应局部区域大小
                C=8  # 常数，用于调整阈值
            )
            
            # 4. 降噪处理
            # 使用形态学操作去除小噪点
            kernel = np.ones((2,2), np.uint8)
            img_array = cv2.morphologyEx(img_array, cv2.MORPH_OPEN, kernel)
            
            # 转回PIL格式
            image = Image.fromarray(img_array)
            
            # 5. 图像锐化
            enhancer = ImageEnhance.Sharpness(image)
            image = enhancer.enhance(2.5)  # 增加锐化程度
            
            # 6. 调整图像大小
            # 对于小字体，增加放大倍数
            scale_factor = 3.0  # 增加放大倍数
            new_width = int(image.width * scale_factor)
            new_height = int(image.height * scale_factor)
            image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            # 7. 倾斜校正
            try:
                # 转换为OpenCV格式
                img_cv = np.array(image)
                
                # 检测边缘
                edges = cv2.Canny(img_cv, 50, 150, apertureSize=3)
                
                # 霍夫变换检测直线
                lines = cv2.HoughLines(edges, 1, np.pi/180, threshold=100)
                
                if lines is not None:
                    angles = []
                    for rho, theta in lines[0]:
                        angle = theta * 180 / np.pi
                        if angle < 45:
                            angles.append(angle)
                        elif angle > 135:
                            angles.append(angle - 180)
                    
                    if angles:
                        median_angle = np.median(angles)
                        if abs(median_angle) > 0.5:
                            image = image.rotate(-median_angle, Image.Resampling.BICUBIC, expand=True)
            
            except Exception as e:
                print(f"倾斜校正失败: {str(e)}")
            
            # 8. 添加边距
            padding = 100  # 增加边距
            new_image = Image.new('L', (image.width + 2*padding, image.height + 2*padding), 255)
            new_image.paste(image, (padding, padding))
            image = new_image
            
            return image
            
        except Exception as e:
            print(f"图像预处理时出错: {str(e)}")
            return image  # 如果处理失败，返回原始图像

def main():
    """主函数：处理用户输入并执行转换流程"""
    try:
        converter = PDFToMarkdown()
        
        # 获取用户输入的文件名
        input_name = input('请输入要转换的文件名(不带后缀): ')
        
        # 选择清理级别
        print("\n请选择文本清理级别：")
        print("0 - 不清理（保留所有OCR识别内容）")
        print("1 - 适当清理（仅清理确定的无用内容）")
        print("2 - 强化清理（更积极地清理可能的干扰内容）")
        
        clean_level = -1
        while clean_level not in [0, 1, 2]:
            try:
                clean_level = int(input("请输入清理级别(0-2): "))
            except ValueError:
                print("请输入有效的数字(0-2)")
        
        # 检查所有支持的文件格式
        found_file = None
        for category in SUPPORTED_FORMATS.values():
            for ext in category:
                test_path = f'{input_name}{ext}'
                if check_file_exists(test_path):
                    found_file = test_path
                    break
            if found_file:
                break
        
        # 检查文本格式
        if not found_file:
            for ext in SUPPORTED_TEXT_FORMATS:
                test_path = f'{input_name}{ext}'
                if check_file_exists(test_path):
                    found_file = test_path
                    break
        
        if not found_file:
            print(f"错误：未找到文件 {input_name}.*")
            print("支持的文件格式：")
            for category, extensions in SUPPORTED_FORMATS.items():
                print(f"- {category.title()}文件:", ", ".join(extensions))
            print("- 文本文件:", ", ".join(SUPPORTED_TEXT_FORMATS))
            return
        
        # 执行转换
        output_path = f'{input_name}.md'
        
        # 修改 process_file 方法调用，传入清理级别
        result = converter.process_file(found_file, output_path, clean_level)
        
        if result:
            print(f'转换完成，结果保存在：{result["original"]}')
            if result['translated']:
                print(f'翻译后的中文版本保存在：{result["translated_path"]}')
        else:
            print('转换失败')
            
    except KeyboardInterrupt:
        print("\n程序被用户中断")
    except Exception as e:
        print(f"程序运行出错: {str(e)}")

if __name__ == '__main__':
    main()